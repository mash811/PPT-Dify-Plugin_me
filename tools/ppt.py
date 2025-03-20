from collections.abc import Generator
import os
import tempfile
import io
import re
import markdown
from typing import Any, List, Dict, Optional
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import PP_PLACEHOLDER

from dify_plugin import Tool
from dify_plugin.entities.tool import ToolInvokeMessage

class PptTool(Tool):
    def _invoke(self, tool_parameters: dict[str, Any]) -> Generator[ToolInvokeMessage]:
        # Get markdown content from parameters
        md_content = tool_parameters.get("markdown_content", "")
        title = tool_parameters.get("title", "Presentation")
        theme = tool_parameters.get("theme", "default")
        
        if not md_content:
            yield self.create_text_message("No markdown content provided.")
            return
        
        try:
            # Convert markdown to PPTX
            pptx_bytes = self._convert_markdown_to_pptx(md_content, title, theme)
            
            # Create a filename
            filename = f"{title.replace(' ', '_')}.pptx"
            
            # Return success message
            yield self.create_text_message(f"PowerPoint presentation '{title}' generated successfully")
            
            # Return the document data as a blob
            yield self.create_blob_message(
                blob=pptx_bytes, 
                meta={
                    "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "filename": filename
                }
            )
        except Exception as e:
            yield self.create_text_message(f"Error converting markdown to PPTX: {str(e)}")
    
    def _convert_markdown_to_pptx(self, md_content: str, title: str, theme: str = "default") -> bytes:
        """Convert markdown content to PowerPoint presentation"""
        # Create a new presentation with optional template
        prs = self._create_presentation(theme)
        
        # Check if the content has slide separators (---)
        if self._has_slide_separators(md_content):
            # Process markdown with slide separators
            return self._process_with_separators(md_content, title, prs)
        else:
            # Process markdown without separators (using headers as slide dividers)
            # Convert markdown to HTML with extensions
            html_content = markdown.markdown(
                md_content,
                extensions=[
                    'markdown.extensions.tables',
                    'markdown.extensions.fenced_code',
                    'markdown.extensions.codehilite',
                    'markdown.extensions.nl2br',
                    'markdown.extensions.sane_lists'
                ]
            )
            
            # Parse HTML content
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Check if presentation metadata is included
            metadata = self._extract_metadata(md_content)
            
            # Process the HTML elements to create slides
            self._create_slides_from_html(prs, soup, title, metadata)
            
            # Save presentation to a bytes buffer
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            
            return pptx_bytes.getvalue()
    
    def _has_slide_separators(self, md_content: str) -> bool:
        """Check if the markdown content uses slide separators (---)"""
        # Look for patterns like \n---\n which indicate slide separators
        return bool(re.search(r'\n-{3,}\n', md_content))
    
    def _process_with_separators(self, md_content: str, title: str, prs: Presentation) -> bytes:
        """Process markdown content that uses slide separators"""
        # Split content by slide separators
        slide_contents = re.split(r'\n-{3,}\n', md_content)
        
        # Extract metadata from the first slide if any
        metadata = self._extract_metadata(slide_contents[0] if slide_contents else "")
        
        # Create title slide from the first content block
        first_slide_content = slide_contents[0] if slide_contents else ""
        html_content = markdown.markdown(
            first_slide_content,
            extensions=[
                'markdown.extensions.tables',
                'markdown.extensions.fenced_code',
                'markdown.extensions.codehilite',
                'markdown.extensions.nl2br',
                'markdown.extensions.sane_lists'
            ]
        )
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Find main title (h1) if any
        main_title = title
        h1_tag = soup.find('h1')
        if h1_tag:
            main_title = h1_tag.get_text().strip()
            # Remove h1 from content since it will be in the title
            h1_tag.decompose()
        
        # Create title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        # Set the title
        title_shape = title_slide.shapes.title
        title_shape.text = main_title
        
        # Set subtitle if available
        if len(title_slide.placeholders) > 1:
            try:
                subtitle = title_slide.placeholders[1]
                h2_tag = soup.find('h2')
                if h2_tag:
                    subtitle.text = h2_tag.get_text().strip()
                    # Remove h2 from content since it will be in subtitle
                    h2_tag.decompose()
                else:
                    author = metadata.get("author", "")
                    date = metadata.get("date", "")
                    
                    if author and date:
                        subtitle.text = f"{author} | {date}"
                    elif author:
                        subtitle.text = author
                    elif date:
                        subtitle.text = date
            except (IndexError, KeyError):
                pass
                
        # Process any remaining content in the first slide (like bullet points)
        remaining_content = list(soup.children)
        if remaining_content:
            # Find a suitable body placeholder in title slide
            body_shape = self._find_body_shape(title_slide)
            if body_shape:
                self._add_content_to_slide(body_shape, remaining_content)
        
        # Process remaining slides
        for slide_content in slide_contents[1:]:
            # Skip empty slides
            if not slide_content.strip():
                continue
                
            # Convert slide content to HTML
            html_content = markdown.markdown(
                slide_content,
                extensions=[
                    'markdown.extensions.tables',
                    'markdown.extensions.fenced_code',
                    'markdown.extensions.codehilite',
                    'markdown.extensions.nl2br',
                    'markdown.extensions.sane_lists'
                ]
            )
            
            # Parse HTML
            slide_soup = BeautifulSoup(html_content, 'html.parser')
            
            # Determine slide layout and title
            slide_title = ""
            h1_tag = slide_soup.find('h1')
            h2_tag = slide_soup.find('h2')
            
            if h1_tag:
                slide_title = h1_tag.get_text().strip()
                h1_tag.decompose()  # Remove from content as it will be the slide title
            elif h2_tag:
                slide_title = h2_tag.get_text().strip()
                h2_tag.decompose()  # Remove from content as it will be the slide title
            
            # Determine slide layout
            has_table = bool(slide_soup.find('table'))
            if has_table:
                try:
                    slide_layout = prs.slide_layouts[3]  # Title and Content with Table
                except IndexError:
                    slide_layout = prs.slide_layouts[1]  # Default to Title and Content
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content
            
            # Create slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if slide.shapes.title:
                slide.shapes.title.text = slide_title
            
            # Get all remaining content elements
            content_elements = list(slide_soup.children)
            
            # Add content if there are any elements left
            if content_elements:
                try:
                    # Find a suitable content placeholder - first try body placeholder
                    body_shape = self._find_body_shape(slide)
                    
                    if body_shape:
                        self._add_content_to_slide(body_shape, content_elements)
                    else:
                        # If no body placeholder found, create a textbox
                        left = Inches(1)
                        top = Inches(2)
                        width = Inches(8)
                        height = Inches(4)
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        self._add_content_to_slide(textbox, content_elements)
                        
                except Exception as e:
                    # If there's a problem, create a textbox and try again
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4)
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    self._add_content_to_slide(textbox, content_elements)
        
        # Save presentation to a bytes buffer
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        return pptx_bytes.getvalue()
    
    def _find_body_shape(self, slide):
        """Find the body shape in a slide (content placeholder)"""
        # First try to find a placeholder with type BODY
        body_shape = self._get_placeholder(slide, PP_PLACEHOLDER.BODY)
        if body_shape:
            return body_shape
            
        # If that fails, check for content placeholder (usually index 1)
        if len(slide.placeholders) > 1:
            try:
                body_candidate = slide.placeholders[1]
                if body_candidate.has_text_frame:
                    return body_candidate
            except (IndexError, KeyError):
                pass
                
        # As a last resort, look for any placeholder with a text frame that isn't the title
        for shape in slide.placeholders:
            if shape != slide.shapes.title and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                return shape
                
        # If all else fails, look at all shapes
        for shape in slide.shapes:
            if shape != slide.shapes.title and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                return shape
                
        return None
    
    def _create_presentation(self, theme: str) -> Presentation:
        """Create a new presentation with an optional template"""
        template_path = None
        
        # Check if a template exists for the given theme
        if theme != "default":
            # Look for user-specified template
            template_name = f"{theme}.pptx"
            possible_paths = [
                os.path.join(os.path.dirname(__file__), "..", "_assets", "templates", template_name),
                os.path.join(os.path.dirname(__file__), "..", template_name)
            ]
            
            for path in possible_paths:
                if os.path.exists(path):
                    template_path = path
                    break
        
        # If a template path was found, use it
        if template_path:
            return Presentation(template_path)
        else:
            # Use default blank presentation
            return Presentation()
    
    def _extract_metadata(self, md_content: str) -> Dict[str, Any]:
        """Extract metadata from markdown content (if any)"""
        metadata = {}
        
        # Check for metadata at the start of the document
        lines = md_content.split('\n')
        in_metadata = True
        metadata_lines = []
        
        for line in lines:
            if line.strip() == "" and in_metadata:
                in_metadata = False
                break
            
            if in_metadata and ":" in line:
                key, value = line.split(":", 1)
                metadata[key.strip().lower()] = value.strip()
                metadata_lines.append(line)
        
        return metadata
    
    def _create_slides_from_html(self, prs: Presentation, soup: BeautifulSoup, presentation_title: str, metadata: Dict[str, Any] = None):
        """Create slides from HTML content"""
        if metadata is None:
            metadata = {}
            
        # First create a title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        
        # Set the title
        title_shape = title_slide.shapes.title
        title_shape.text = presentation_title
        
        # Try to set subtitle if the slide has a subtitle placeholder
        if len(title_slide.placeholders) > 1:
            try:
                subtitle = title_slide.placeholders[1]  # Index 1 is typically the subtitle
                author = metadata.get("author", "")
                date = metadata.get("date", "")
                
                if author and date:
                    subtitle.text = f"{author} | {date}"
                elif author:
                    subtitle.text = author
                elif date:
                    subtitle.text = date
            except (IndexError, KeyError):
                pass
        
        # Find all tags that could be headings or content
        all_elements = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'pre', 'code', 'table'])
        
        # Group elements by slide (h1/h2 tags create new slides)
        slides_content = []
        current_group = []
        
        for elem in all_elements:
            if elem.name in ['h1', 'h2']:
                # If we have a previous slide group, save it
                if current_group:
                    slides_content.append(current_group)
                # Start a new slide group with this heading
                current_group = [elem]
            elif current_group:  # Add to current slide group if one exists
                current_group.append(elem)
            # If no current group exists yet, create one if this is a heading
            elif elem.name in ['h3', 'h4', 'h5', 'h6']:
                current_group = [elem]
        
        # Add the last slide group if it exists
        if current_group:
            slides_content.append(current_group)
        
        # Create slides from the grouped content
        for slide_elements in slides_content:
            # Get the first element as heading
            heading = slide_elements[0]
            content = slide_elements[1:]
            
            # Determine slide layout based on heading level
            if heading.name == 'h1':
                slide_layout = prs.slide_layouts[0]  # Title slide
            elif any(elem.name == 'table' for elem in content):
                # If content contains a table, use a slide with title and content
                try:
                    slide_layout = prs.slide_layouts[3]  # Title and Content
                except IndexError:
                    slide_layout = prs.slide_layouts[1]  # Default to Title and Content
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content slide
            
            # Create slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = heading.get_text()
            
            # Add content, but only if there is some
            if content:
                try:
                    # Find a suitable content placeholder
                    body_shape = self._find_body_shape(slide)
                    
                    if body_shape:
                        self._add_content_to_slide(body_shape, content)
                    else:
                        # If no body placeholder found, create a textbox
                        left = Inches(1)
                        top = Inches(2)
                        width = Inches(8)
                        height = Inches(4)
                        textbox = slide.shapes.add_textbox(left, top, width, height)
                        self._add_content_to_slide(textbox, content)
                        
                except Exception as e:
                    # If there's a problem, create a textbox and try again
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4)
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    self._add_content_to_slide(textbox, content)
    
    def _get_placeholder(self, slide, placeholder_type):
        """Get a placeholder by type"""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == placeholder_type:
                return shape
        return None
    
    def _add_content_to_slide(self, shape, content_elements):
        """Add HTML content elements to a slide shape"""
        text_frame = shape.text_frame
        text_frame.clear()  # Clear any existing text
        text_frame.word_wrap = True
        
        # Process each content element
        for element in content_elements:
            # Skip if it's just whitespace
            if isinstance(element, NavigableString) and not element.strip():
                continue
                
            # Handle different types of elements
            if element.name == 'p':
                p = text_frame.add_paragraph()
                # Get text and remove extra whitespace
                text_content = element.get_text().strip()
                if text_content:
                    p.text = text_content
            
            elif element.name in ['ul', 'ol']:
                self._add_list_to_textframe(text_frame, element, is_ordered=element.name == 'ol')
            
            elif element.name == 'table':
                # Convert table to text representation in PowerPoint
                self._add_table_as_text(text_frame, element)
            
            elif element.name in ['pre', 'code']:
                # Code blocks
                p = text_frame.add_paragraph()
                code_text = element.get_text().strip()
                if code_text:
                    p.text = code_text
                    # Format as code (monospace font)
                    for run in p.runs:
                        run.font.name = 'Courier New'
                        run.font.size = Pt(10)
            
            elif element.name in ['h3', 'h4', 'h5', 'h6']:
                # Add subheadings within the slide with appropriate formatting
                p = text_frame.add_paragraph()
                text_content = element.get_text().strip()
                if text_content:
                    p.text = text_content
                    # Format as subheading
                    for run in p.runs:
                        run.bold = True
                        # Set size based on heading level
                        if element.name == 'h3':
                            run.font.size = Pt(18)
                        elif element.name == 'h4':
                            run.font.size = Pt(16)
                        elif element.name == 'h5':
                            run.font.size = Pt(14)
                        elif element.name == 'h6':
                            run.font.size = Pt(12)
    
    def _add_list_to_textframe(self, text_frame, list_element, is_ordered=False, level=0):
        """Add a list to a text frame with proper indentation"""
        # Get all list items at this level
        list_items = list_element.find_all('li', recursive=False)
        
        for i, item in enumerate(list_items):
            # Create a new paragraph for this list item
            p = text_frame.add_paragraph()
            
            # Extract text directly from this li element (not including nested lists)
            text_parts = []
            for child in item.children:
                if isinstance(child, NavigableString):
                    text_parts.append(child.strip())
                elif child.name not in ['ul', 'ol']:  # Don't include nested lists
                    text_parts.append(child.get_text().strip())
            
            # Join text parts and set paragraph text
            item_text = ' '.join(text_parts).strip()
            
            # For ordered lists, manually add the number at the beginning
            if is_ordered:
                p.text = f"{i+1}. {item_text}"
            else:
                # For unordered lists, we'll use PowerPoint's built-in bullet formatting
                p.text = item_text
            
            # Set the indentation level
            p.level = level
            
            # Explicitly set paragraph as bullet point using direct XML manipulation
            # This ensures bullets appear correctly regardless of PPT template
            self._ensure_bullet_formatting(p, is_ordered, i+1)
            
            # Process any nested lists
            nested_lists = item.find_all(['ul', 'ol'], recursive=False)
            for nested_list in nested_lists:
                self._add_list_to_textframe(
                    text_frame, 
                    nested_list, 
                    is_ordered=nested_list.name == 'ol',
                    level=level+1
                )
    
    def _ensure_bullet_formatting(self, paragraph, is_ordered=False, number=1):
        """Ensure bullet or numbering is applied to paragraph using direct XML approach"""
        # We'll try multiple approaches in order of preference
        try:
            # First try the simplest approach: set the bullet property
            self._set_bullet_property(paragraph, is_ordered)
            return
        except Exception as e:
            pass
            
        try:
            # Next try using the parse_xml approach
            self._apply_bullet_via_parse_xml(paragraph, is_ordered)
            return
        except Exception as e:
            pass
            
        try:
            # Try direct XML manipulation if the above methods fail
            self._apply_bullet_via_direct_xml(paragraph, is_ordered)
            return
        except Exception as e:
            pass
            
        try:
            # Final attempt: use raw XML string approach
            self._apply_bullet_via_xml_string(paragraph, is_ordered)
        except Exception as e:
            # If all methods fail, we'll have to accept that bullets might not display
            pass
    
    def _set_bullet_property(self, paragraph, is_ordered=False):
        """Try to use the built-in python-pptx bullet property"""
        try:
            if hasattr(paragraph, 'bullet') and paragraph.bullet is not None:
                paragraph.bullet.visible = True
                if is_ordered:
                    paragraph.bullet.number_style = 'arabicPeriod'
                else:
                    paragraph.bullet.character = '•'
                return True
        except:
            return False
    
    def _apply_bullet_via_xml_string(self, paragraph, is_ordered=False):
        """A last-resort method to set bullet formatting using string XML"""
        # Get the internal lxml element
        p_element = paragraph._element
        
        # Define the XML string for bullet points
        if is_ordered:
            xml_string = '<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="arabicPeriod" startAt="1"/>'
        else:
            xml_string = '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>'
        
        # Define the pPr element
        pPr_string = '<a:pPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">' + xml_string + '</a:pPr>'
        
        # Import the lxml module within the function to handle the XML
        from lxml import etree
        
        # Remove any existing pPr element
        for child in list(p_element):
            if child.tag.endswith('pPr'):
                p_element.remove(child)
        
        # Create and add the new pPr element with bullet formatting
        parser = etree.XMLParser(remove_blank_text=True)
        new_pPr = etree.fromstring(pPr_string, parser)
        p_element.insert(0, new_pPr)
    
    def _apply_bullet_via_parse_xml(self, paragraph, is_ordered=False):
        """Apply bullet or numbering via parse_xml method"""
        from pptx.oxml import parse_xml
        
        # Get paragraph properties element
        pPr = paragraph._p.get_or_add_pPr()
        
        # First, remove any existing bullet properties that might conflict
        for tag in ['{*}buNone', '{*}buChar', '{*}buAutoNum']:
            for element in pPr.findall(f'.//{tag}'):
                if element.getparent() is not None:
                    element.getparent().remove(element)
        
        if is_ordered:
            # Add auto-numbering for ordered lists
            buAutoNum = parse_xml(f'<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="arabicPeriod" startAt="1"/>')
            pPr.append(buAutoNum)
        else:
            # Add bullet character for unordered lists
            buChar = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>')
            pPr.append(buChar)
    
    def _apply_bullet_via_direct_xml(self, paragraph, is_ordered=False):
        """Apply bullet via direct element creation (fallback method)"""
        from lxml import etree
        
        # Define namespace map
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
        }
        
        # Get paragraph element
        p_element = paragraph._p
        
        # Find existing pPr or create it
        pPr = p_element.find('.//a:pPr', namespaces=nsmap)
        if pPr is None:
            # Need to create a new pPr element
            if len(p_element) > 0 and p_element[0].tag.endswith('pPr'):
                # Use existing one if it's there
                pPr = p_element[0]
            else:
                # Create a new one - this is tricky with namespaces, so we'll use a tag with namespace
                pPr = etree.Element('{{{0}}}pPr'.format(nsmap['a']))
                # Insert at beginning
                if len(p_element) > 0:
                    p_element.insert(0, pPr)
                else:
                    p_element.append(pPr)
        
        # Remove any existing bullet elements using full namespace tags
        for child in list(pPr):
            tag = child.tag
            if tag.endswith('buNone') or tag.endswith('buChar') or tag.endswith('buAutoNum'):
                pPr.remove(child)
        
        # Create the appropriate bullet element with proper namespace
        if is_ordered:
            # For ordered list
            buAutoNum = etree.Element('{{{0}}}buAutoNum'.format(nsmap['a']))
            buAutoNum.set('type', 'arabicPeriod')
            buAutoNum.set('startAt', '1')
            pPr.append(buAutoNum)
        else:
            # For bullet list
            buChar = etree.Element('{{{0}}}buChar'.format(nsmap['a']))
            buChar.set('char', '•')
            pPr.append(buChar)
        
        # Try setting bullet font as well 
        buFont = etree.Element('{{{0}}}buFont'.format(nsmap['a']))
        buFont.set('typeface', '+mj-lt')  # Default bullet font
        pPr.append(buFont)
    
    def _add_table_as_text(self, text_frame, table_element):
        """Convert HTML table to text representation"""
        # Get table rows
        rows = table_element.find_all('tr')
        
        # Process headers if present
        headers = table_element.find_all('th')
        if headers:
            header_text = ' | '.join([header.get_text().strip() for header in headers])
            p = text_frame.add_paragraph()
            p.text = header_text
            
            # Add separator line
            separator = '-' * len(header_text)
            p = text_frame.add_paragraph()
            p.text = separator
        
        # Process rows
        for row in rows:
            cells = row.find_all(['td', 'th'])
            if cells:
                p = text_frame.add_paragraph()
                p.text = ' | '.join([cell.get_text().strip() for cell in cells])
